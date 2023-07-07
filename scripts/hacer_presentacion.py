# -*- coding: utf-8 -*-
"""
Created on Thu Jul  6 08:58:20 2023

@author: AARON.RAMIREZ
"""
import os
from pptx import Presentation
from pptx.util import Inches

equipos_nom = { 
    'Operación Estratégica': ['LILIANA AVILA LOPEZ',
      'VERONICA ITZEL JIMENEZ GONZALEZ',#no estan en base de datos
      'MARIANA RIOS MARTINEZ',
      'VICTOR AYAX QUINTERO ORTEGA',
      'DANIEL LOPEZ SANCHEZ'],#no estan en base de datos
    'Integración de Información': ['NALLELY BECERRIL DAVILA',
      'ROGELIO ROSALES MORALES',
      'JOSE MANUEL OCOMATL OLAYA',
      'HUGO GONZALEZ VALDEZ',
      'RUBI SAMANTHA MEDRANO MARTINEZ',
      'JOSE ANTONIO CHAVEZ CASTILLO',
      'VICTOR RAMIRO ESPINA CASAS',
      'ANA AGLAE FLORES AGUILAR'],
    'Control y Logística': ['ALEXEI PRADEL HERNANDEZ',
      # 'MA. GUADALUPE ADAME SALGADO',
      'KARLA FABIOLA ACEVEDO BERNARDINO',
      # 'ANTONIO ROMERO LEYVA',
      'YOLISMA IVETTE LOPEZ CERON',
      # 'DIANA LETICIA ALCALA GONZALEZ',
      'EDWIN HECTOR PINEDA LOPEZ']
    }
siglas = {'Operación Estratégica':'DOE',
          'Integración de Información': 'DII',
          'Control y Logística':'DCL'
          }
def gen_presentacion():
    """
    Al ser llamada, la función busca en la carpeta de utilidades todos
    los archivos para ordenarlos y pegarlos en una presentación con
    el desempeño de los responsables revisores. 
    """
    prs = Presentation('Utilidades_app\Plantilla_institucional.pptx')
    #busca todos los archivos de la carpeta utilidades
    archivos = os.listdir('Utilidades_app')
    for equipo in equipos_nom:
        p = prs.slides.add_slide(prs.slide_layouts[5])
        titulo = p.shapes.title 
        titulo.text = equipo
        for miembro in equipos_nom[equipo]:
            for archivo in archivos:
                if miembro in archivo:
                   q = prs.slides.add_slide(prs.slide_layouts[7])
                   titulo = q.shapes.title
                   titulo.text = siglas[equipo]+'-Responsable revisor: '+miembro 
                   foto = q.shapes.add_picture(f'Utilidades_app\{miembro}.png',
                                               Inches(0.1),
                                               Inches(0.1))
    prs.save('informe_responsables_revisores.pptx')

def borrar_imagenes():
    for equipo in equipos_nom:
        for miembro in equipos_nom[equipo]:
            # print('reviso ',miembro)
            if os.path.exists(f'Utilidades_app/{miembro}.png'):
                # print('borro ',miembro)
                os.remove(f'Utilidades_app/{miembro}.png')         
